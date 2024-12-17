import streamlit as st
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import os

# Streamlit App Title
st.title("PDF를 이미지로 변환하여 PPTX로 저장하기 📄➡️📽️")
st.write("PDF 파일을 업로드하면 각 페이지를 이미지로 변환하고, 이를 PPTX 파일로 만들어 다운로드할 수 있어요! 😊")

# PDF 파일 업로드
uploaded_pdf = st.file_uploader("PDF 파일을 업로드 해주세요.", type=["pdf"])

if uploaded_pdf is not None:
    # PDF를 BytesIO로 읽기
    pdf_bytes = uploaded_pdf.read()
    st.info("PDF 파일을 이미지로 변환하는 중입니다. 잠시만 기다려주세요... 🕒")
    
    # PDF를 이미지로 변환
    images = convert_from_bytes(pdf_bytes)
    st.success(f"PDF에서 {len(images)} 페이지를 이미지로 변환했어요! 🎉")
    
    # PPTX 파일 생성
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # 빈 슬라이드 레이아웃
    
    for i, image in enumerate(images):
        slide = prs.slides.add_slide(blank_slide_layout)
        image_stream = BytesIO()
        image.save(image_stream, format="PNG")
        image_size = image.size  # 이미지 사이즈 (width, height)
        width_in_inches = image_size[0] / 96  # 이미지 너비를 인치로 변환 (96 dpi 기준)
        height_in_inches = image_size[1] / 96  # 이미지 높이를 인치로 변환
        
        # 이미지 삽입
        slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=Inches(width_in_inches), height=Inches(height_in_inches))
        st.write(f"페이지 {i+1} 변환 완료! 🖼️")
    
    # PPTX 파일 저장
    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    # 다운로드 버튼 제공
    st.success("모든 페이지를 PPTX 파일로 변환했어요! 🎯")
    st.download_button(
        label="PPTX 파일 다운로드 ⬇️",
        data=pptx_bytes,
        file_name="converted_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
else:
    st.warning("PDF 파일을 업로드해주세요! ☝️")

st.write("이 도구는 PDF의 각 페이지를 이미지로 변환하여 PPTX 파일로 저장해줍니다. 변환된 PPTX 파일은 페이지 크기에 맞게 조정됩니다! 🚀")
